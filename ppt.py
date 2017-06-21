#!/usr/bin/python

from pptx import Presentation
from pptx.util import Inches
x=1
prs = Presentation()
while(x<34):#change 34 by your page value ex: 33 pages 33+1 or x<=33
	img_path = str(x)+'.jpg'
	
	blank_slide_layout = prs.slide_layouts[6]
	slide = prs.slides.add_slide(blank_slide_layout)
	left = top = Inches(1)
	width = Inches(10.0)
	height = Inches(7.50)
	pic = slide.shapes.add_picture(img_path, left, top,width,height)
	x=x+1


prs.save('test.pptx')
